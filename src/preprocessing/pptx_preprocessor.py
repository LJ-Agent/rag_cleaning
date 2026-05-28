"""PowerPoint (PPTX) preprocessor — extract slides as pages, text/shapes/tables/images per slide."""

import io
from uuid import uuid4

from common.config_loader import get_config
from common.exception.exceptions import PreprocessingException
from common.models.document import (
    BoundingBox,
    Document,
    DocumentMetadata,
    ElementRole,
    ImageElement,
    Page,
    TableCell,
    TableElement,
    TextElement,
)
from common.util.logger import bind_trace_id, get_logger
from common.util.utils import md5_bytes
from preprocessing.base import BasePreprocessor

logger = get_logger()


class PptxPreprocessor(BasePreprocessor):
    """Extract raw elements from PowerPoint presentations. Each slide -> one page."""

    format_type = "pptx"
    supported_extensions = {"pptx", "ppt"}
    supported_mime_types = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }

    def __init__(self):
        cfg = get_config()["preprocessing"]["pptx"]
        self._extract_notes = cfg.get("extract_notes", True)
        self._extract_master = cfg.get("extract_master_slide", False)

    def extract(self, file_data: bytes, file_name: str) -> Document:
        log = bind_trace_id(str(uuid4())[:8])
        doc_id = md5_bytes(file_data)[:16]
        doc = Document(doc_id=doc_id)

        doc.metadata = DocumentMetadata(
            source_format="pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_size_bytes=len(file_data),
            file_md5=md5_bytes(file_data),
        )
        doc.log_stage("pptx_preprocess_start")

        try:
            from pptx import Presentation
            from pptx.shapes.base import BaseShape
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(io.BytesIO(file_data))

            # Metadata
            if prs.core_properties:
                doc.metadata.title = prs.core_properties.title or ""
                doc.metadata.author = prs.core_properties.author or ""
                doc.metadata.created_at = str(prs.core_properties.created) if prs.core_properties.created else ""

            for slide_idx, slide in enumerate(prs.slides):
                page = Page(page_number=slide_idx + 1)
                elem_idx = 0
                full_text = ""

                # Extract shapes in z-order
                for shape in slide.shapes:
                    shape_type = shape.shape_type

                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            full_text += text + "\n"

                            role = ElementRole.PARAGRAPH
                            if shape.is_placeholder:
                                ph = shape.placeholder_format
                                if ph and ph.type is not None:
                                    from pptx.enum.shapes import PP_PLACEHOLDER
                                    if ph.type == PP_PLACEHOLDER.TITLE:
                                        role = ElementRole.HEADING

                            elem = TextElement(
                                element_id=self._make_element_id(doc_id, slide_idx + 1, elem_idx, "txt"),
                                role=role,
                                page_numbers=[slide_idx + 1],
                                text=text,
                            )
                            page.elements.append(elem)
                            elem_idx += 1

                    elif shape.has_table:
                        table_data = shape.table
                        rows = []
                        for row in table_data.rows:
                            cells = [TableCell(text=cell.text.strip()) for cell in row.cells]
                            rows.append(cells)

                        table = TableElement(
                            element_id=self._make_element_id(doc_id, slide_idx + 1, elem_idx, "tbl"),
                            role=ElementRole.TABLE,
                            page_numbers=[slide_idx + 1],
                            rows=rows,
                            column_count=len(table_data.columns),
                            row_count=len(table_data.rows),
                        )
                        page.elements.append(table)
                        elem_idx += 1
                        doc.metadata.has_tables = True

                    elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_data = shape.image.blob
                            img_hash = md5_bytes(img_data)
                            img = ImageElement(
                                element_id=self._make_element_id(doc_id, slide_idx + 1, elem_idx, "img"),
                                role=ElementRole.IMAGE,
                                page_numbers=[slide_idx + 1],
                                image_data=img_data,
                                image_hash=img_hash,
                                width=int(shape.width),
                                height=int(shape.height),
                                format=shape.image.content_type.split("/")[-1] if shape.image.content_type else "png",
                            )
                            page.elements.append(img)
                            elem_idx += 1
                            doc.metadata.has_images = True
                        except Exception:
                            pass

                # Extract speaker notes
                if self._extract_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        note = TextElement(
                            element_id=self._make_element_id(doc_id, slide_idx + 1, elem_idx, "notes"),
                            role=ElementRole.QUOTE,
                            page_numbers=[slide_idx + 1],
                            text=f"[Speaker Notes]: {notes_text}",
                        )
                        page.elements.append(note)
                        full_text += notes_text + "\n"
                        elem_idx += 1

                page.text_content = full_text
                doc.metadata.char_count += len(full_text)
                doc.metadata.word_count += len(full_text.split())
                doc.pages.append(page)

            doc.metadata.page_count = len(doc.pages)

        except Exception as e:
            raise PreprocessingException(f"PPTX extraction failed: {e}", format_type="pptx")

        doc.log_stage("pptx_preprocess_done")
        return doc
